import asyncio
import json
import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from unittest.mock import AsyncMock, patch

from fastapi import UploadFile
from pptx import Presentation

from backend.services import uploads
from backend.services.deck_review import review_deck_session
from backend.services.model_router import model_supports_vision, provider_catalog, stream_chat_completion


class DeckReviewTests(unittest.TestCase):
    def setUp(self) -> None:
        self.tempdir = tempfile.TemporaryDirectory()
        self.original_data_dir = uploads.DATA_DIR
        self.original_uploads_dir = uploads.UPLOADS_DIR
        self.original_max_upload_bytes = uploads.MAX_UPLOAD_BYTES
        self.original_max_uploads = uploads.MAX_UPLOADS_PER_SESSION
        uploads.DATA_DIR = Path(self.tempdir.name)
        uploads.UPLOADS_DIR = uploads.DATA_DIR / "session_uploads"

    def tearDown(self) -> None:
        uploads.DATA_DIR = self.original_data_dir
        uploads.UPLOADS_DIR = self.original_uploads_dir
        uploads.MAX_UPLOAD_BYTES = self.original_max_upload_bytes
        uploads.MAX_UPLOADS_PER_SESSION = self.original_max_uploads
        self.tempdir.cleanup()

    def _pptx_bytes(self) -> bytes:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Risk intelligence for hedge funds"
        slide.placeholders[1].text = "Problem: post-trade anomalies are caught too late.\nCustomer discovery: 12 interviews."
        slide_two = prs.slides.add_slide(prs.slide_layouts[1])
        slide_two.shapes.title.text = "Market and implementation"
        slide_two.placeholders[1].text = "TAM / SAM / SOM\nRoadmap: pilot in semester one, 18-month rollout."
        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def test_model_supports_vision_heuristics(self) -> None:
        self.assertTrue(model_supports_vision("ollama", "qwen2.5vl:7b"))
        self.assertTrue(model_supports_vision("openai", "gpt-4o"))
        self.assertTrue(model_supports_vision("openai", "gpt-4.1"))
        self.assertFalse(model_supports_vision("ollama", "qwen3:8b"))
        self.assertFalse(model_supports_vision("cerebras", "qwen-3-32b"))

    def test_provider_catalog_marks_server_configured_keys(self) -> None:
        with patch.dict("os.environ", {"GROQ_API_KEY": "test-key"}, clear=False):
            groq = next(item for item in provider_catalog() if item["key"] == "groq")
        self.assertTrue(groq["serverConfigured"])
        self.assertTrue(groq["defaultSpeedModel"])
        self.assertTrue(groq["openWeight"])

    def test_ingest_upload_builds_pptx_deck_artifact(self) -> None:
        upload = UploadFile(filename="deck.pptx", file=BytesIO(self._pptx_bytes()))
        entry = asyncio.run(uploads.ingest_upload("deck-session", upload))
        artifact = uploads.load_deck_artifact("deck-session")

        self.assertEqual(entry["docType"], "pitch deck")
        self.assertEqual(entry["slideCount"], 2)
        self.assertIsNotNone(artifact)
        self.assertEqual(artifact["format"], "pptx")
        self.assertEqual(artifact["slideCount"], 2)
        self.assertFalse(artifact["hasRenderableSlides"])
        self.assertEqual(artifact["slides"][0]["label"], "Slide 1")

    def test_upload_rejects_unsupported_file_type(self) -> None:
        upload = UploadFile(filename="payload.exe", file=BytesIO(b"not allowed"))

        with self.assertRaisesRegex(ValueError, "Unsupported file type"):
            asyncio.run(uploads.ingest_upload("bad-type", upload))

    def test_upload_rejects_files_over_beta_limit(self) -> None:
        uploads.MAX_UPLOAD_BYTES = 4
        upload = UploadFile(filename="notes.txt", file=BytesIO(b"too large"))

        with self.assertRaisesRegex(ValueError, "File is too large"):
            asyncio.run(uploads.ingest_upload("too-large", upload))

    def test_upload_rejects_too_many_files_in_session(self) -> None:
        uploads.MAX_UPLOADS_PER_SESSION = 1
        first = UploadFile(filename="one.txt", file=BytesIO(b"first note"))
        second = UploadFile(filename="two.txt", file=BytesIO(b"second note"))

        asyncio.run(uploads.ingest_upload("upload-cap", first))
        with self.assertRaisesRegex(ValueError, "already has 1 uploads"):
            asyncio.run(uploads.ingest_upload("upload-cap", second))

    def test_review_deck_session_uses_structured_artifact(self) -> None:
        upload = UploadFile(filename="deck.pptx", file=BytesIO(self._pptx_bytes()))
        asyncio.run(uploads.ingest_upload("deck-review", upload))

        fake_payload = {
            "verdict": "Promising, but still under-proven.",
            "summary": "The deck has a credible problem and market start, but discovery and proof are still thin.",
            "whatWorks": ["Slide 1: the problem statement is concrete enough to understand the pain."],
            "weakPoints": ["Not shown: the competition story is still incomplete."],
            "unprovenClaims": ["Deck-level: the value proposition needs more proof than the current deck shows."],
            "storyFlow": "The flow is logical, but it jumps from problem to market without enough proof in between.",
            "focusedAssessments": [
                {"key": "customer_discovery", "label": "Customer discovery", "status": "partial", "assessment": "Slide 1 mentions interviews, but the learning is still thin.", "refs": ["Slide 1"]},
            ],
            "slideReviews": [
                {"index": 1, "label": "Slide 1", "summary": "Strong problem framing.", "whatWorks": ["Clear pain signal."], "issues": ["Needs more evidence."], "suggestions": ["Add one sharper customer quote or result."], "refs": ["Slide 1"]},
            ],
            "topFixes": ["Add a sharper competition slide.", "Show stronger customer proof."],
            "londonWhaleAssessment": "Not shown in the current deck or context.",
        }

        with patch("backend.services.deck_review.generate_provider_text", new=AsyncMock(return_value={
            "message": json.dumps(fake_payload),
            "model": "qwen3:8b",
            "provider": "ollama",
            "finishReason": "stop",
            "timings": {"firstTokenSeconds": 0.1, "totalSeconds": 0.2},
        })):
            report = asyncio.run(
                review_deck_session(
                    session_id="deck-review",
                    provider="ollama",
                    model="qwen3:8b",
                    user_context="Be tough on proof and story flow.",
                )
            )

        self.assertEqual(report["reviewMode"], "text_transcript")
        self.assertGreaterEqual(report["overallScore"], 0)
        self.assertTrue(report["templateCoverage"])
        self.assertEqual(report["slideReviews"][0]["label"], "Slide 1")
        self.assertTrue(any(item["section"] == "Problem Statement" for item in report["templateCoverage"]))

        title_section = next(item for item in report["templateCoverage"] if item["section"] == "Title Slide")
        self.assertEqual(title_section["refs"], ["Slide 1"])
        self.assertTrue(title_section["evidence"])

        problem_section = next(item for item in report["templateCoverage"] if item["section"] == "Problem Statement")
        self.assertTrue(problem_section["evidence"])
        self.assertIn("what problem is being solved", problem_section["note"].lower())

        feasibility_section = next(item for item in report["templateCoverage"] if item["section"] == "Feasibility")
        self.assertIn("missing", feasibility_section["status"])
        self.assertTrue(feasibility_section["missingItems"])

    def test_stream_chat_completion_continues_after_length_limit(self) -> None:
        async def fake_stream(*args, **kwargs):
            yield "meta", {"responseProfile": "balanced", "model": "test-model", "provider": "openai", "fallbackUsed": False}
            yield "delta", {"delta": "The deck is promising but"}
            yield "complete", {
                "message": "The deck is promising but",
                "responseProfile": "balanced",
                "model": "test-model",
                "provider": "openai",
                "fallbackUsed": False,
                "finishReason": "length",
                "timings": {"firstTokenSeconds": 0.1, "totalSeconds": 0.2},
            }

        async def collect_events():
            items = []
            with patch("backend.services.model_router._stream_from_profile", new=fake_stream), patch(
                "backend.services.model_router.generate_provider_text",
                new=AsyncMock(
                    return_value={
                        "message": "it still needs proof.",
                        "model": "test-model",
                        "provider": "openai",
                        "finishReason": "stop",
                        "timings": {"firstTokenSeconds": 0.1, "totalSeconds": 0.2},
                    }
                ),
            ):
                async for event, payload in stream_chat_completion(
                    system="system",
                    messages=[{"role": "user", "content": "Review the deck"}],
                    response_profile="balanced",
                    provider_override="openai",
                    model_override="gpt-4o",
                    api_key="test-key",
                ):
                    items.append((event, payload))
            return items

        events = asyncio.run(collect_events())
        done_payload = [payload for event, payload in events if event == "complete"][-1]

        self.assertTrue(done_payload["continuedAfterLengthLimit"])
        self.assertEqual(done_payload["continuationCount"], 1)
        self.assertEqual(done_payload["finishReason"], "stop")
        self.assertIn("The deck is promising but it still needs proof.", done_payload["message"])


if __name__ == "__main__":
    unittest.main()

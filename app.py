"""Pitch Deck Mentor — Socratic Thinking Partner for Startup Founders.

Usage:
  python app.py           → User interface  (http://localhost:7860)
  python app.py --admin   → Admin interface (http://localhost:7861)

Model selection (set in .env):
  ANTHROPIC_API_KEY=sk-ant-...   → uses Claude (recommended)
  (no key)                       → falls back to local Ollama / Qwen3:8b
"""

import json
import os
import re
import sys
from pathlib import Path

import gradio as gr

IS_ADMIN = "--admin" in sys.argv

# ── Load .env if present ──────────────────────────────────────────────────────
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# ── RAG + Memory (graceful degradation if deps missing) ───────────────────────
try:
    import rag as _rag
    import memory as _memory
    _memory.init_db()
    RAG_ENABLED = True
except Exception:
    _rag = None   # type: ignore
    _memory = None  # type: ignore
    RAG_ENABLED = False

from prompts import (
    ANALYSIS_SYSTEM_PROMPT,
    build_analysis_prompt,
    build_personalized_opening,
    build_structure_prompt,
    build_system_prompt,
    get_chip_suggestions,
)

# ── Dynamic chip engine ───────────────────────────────────────────────────────
# Maps keyword tuples (found in the mentor's message) → response chip options.
# Chips are short first-person phrases the founder can tap to send as a reply.
_QUESTION_CHIPS: list[tuple[tuple, list[str]]] = [
    (
        ("cac", "customer acquisition cost", "how do you acquire", "acquisition channel", "how are you acquiring customers"),
        ["Our CAC is roughly...", "I haven't measured this yet", "We acquire mainly through...", "Our main channel is..."],
    ),
    (
        ("tam", "total addressable market", "how big is the market", "market size", "addressable market"),
        ["The TAM is roughly...", "I've sized it bottom-up at...", "It's a ₹X crore / $XB market", "I haven't formally sized it yet"],
    ),
    (
        ("ltv", "lifetime value", "customer lifetime", "how long do customers", "how much does a customer"),
        ["Average customer stays X months", "I don't track LTV yet", "A customer spends roughly ₹X total", "LTV to CAC is about X:1"],
    ),
    (
        ("retain", "retention", "churn", "come back", "repeat purchase", "keep them"),
        ["Our retention is X%", "Churn is around X% monthly", "Customers come back every X months", "We haven't measured this formally"],
    ),
    (
        ("revenue", "monetise", "monetize", "pricing", "how do you charge", "willingness to pay", "pay for"),
        ["We charge ₹X per month/unit", "We're pre-revenue right now", "Subscription at ₹X / usage-based", "Revenue is ₹X this month"],
    ),
    (
        ("competitor", "competition", "alternatives", "instead of you", "what do people use today", "do today"),
        ["Main competitors are...", "People currently do it manually", "No direct competitor yet", "We compete mainly on..."],
    ),
    (
        ("team", "co-founder", "who built", "your background", "why you", "unfair advantage"),
        ["I'm the solo founder", "We're a two-person team", "My background is in...", "My co-founder handles..."],
    ),
    (
        ("traction", "users", "customers", "how many", "early signal", "product-market fit", "pmf"),
        ["We have X users / customers", "We're pre-launch", "Growing at X% per month", "Strong word-of-mouth so far"],
    ),
    (
        ("problem", "pain point", "who faces", "how real is", "who has this problem"),
        ["The core problem is...", "The person who faces it is...", "I've experienced this personally", "Here's a concrete example..."],
    ),
    (
        ("why now", "timing", "what changed", "what enabled", "why is this possible now"),
        ["A regulatory / tech shift enabled this", "Behaviour changed post-COVID", "New infrastructure made this possible", "The timing is right because..."],
    ),
    (
        ("raise", "raising", "ask", "how much", "funding round", "investment", "capital"),
        ["We're raising ₹X crore / $X", "Haven't decided the amount yet", "Targeting seed round of...", "We need X months of runway"],
    ),
    (
        ("margin", "gross margin", "unit economics", "contribution margin", "profitable", "profit per"),
        ["Gross margin is X%", "Not profitable at unit level yet", "Each unit earns ₹X after costs", "Payback period is X months"],
    ),
    (
        ("moat", "defensible", "competitive advantage", "what stops", "replicate", "copy"),
        ["Our moat is...", "Network effects protect us", "Proprietary data / IP gives us an edge", "Switching costs are high because..."],
    ),
    (
        ("bottom-up", "build from", "unit of", "price times", "serviceable", "sam"),
        ["Let me build it from the bottom up", "Beachhead is X customers at ₹Y each", "I need help thinking through this", "My SAM estimate is..."],
    ),
    (
        ("insight", "secret", "what do you know", "what do you understand", "that others miss"),
        ["Our insight is...", "We've seen this from inside the industry", "Most people think X, but actually...", "Our access to X is unique"],
    ),
    (
        ("distribution", "go-to-market", "gtm", "reach customers", "sales motion", "how will you sell"),
        ["We sell direct / through partners", "GTM is inbound-led / outbound-led", "Our sales cycle is X weeks", "We're channel-agnostic right now"],
    ),
    (
        ("nps", "net promoter", "love", "satisfaction", "do customers love"),
        ["NPS is around X", "We get strong word-of-mouth", "Customers refer others regularly", "Haven't measured NPS formally yet"],
    ),
]


def generate_contextual_chips(assistant_msg: str, state) -> list[dict]:
    """Generate chips based on keywords in the mentor's last message.

    Falls back to phase/state-based chips when no keyword matches.
    """
    msg_lower = assistant_msg.lower()

    for keywords, chip_texts in _QUESTION_CHIPS:
        if any(kw in msg_lower for kw in keywords):
            result = []
            for i in range(4):
                if i < len(chip_texts):
                    result.append({"text": chip_texts[i], "visible": True})
                else:
                    result.append({"text": "", "visible": False})
            return result

    # No keyword match — fall back to phase/coverage-based chips
    return get_chip_suggestions(state)
from state import ConversationState

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
USE_CLAUDE = bool(ANTHROPIC_API_KEY)

# Fast model: all conversation turns + analysis
CLAUDE_MODEL_FAST    = os.environ.get("CLAUDE_MODEL_FAST",    "claude-haiku-4-5-20251001")
# Premium model: Crystallize output only
CLAUDE_MODEL_PREMIUM = os.environ.get("CLAUDE_MODEL_PREMIUM", "claude-sonnet-4-6")
# Legacy alias kept for .env back-compat
CLAUDE_MODEL = os.environ.get("CLAUDE_MODEL", CLAUDE_MODEL_FAST)

# Local Ollama model (override in .env: OLLAMA_MODEL=gemma3:12b)
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "qwen3:8b")

WELCOME_MESSAGE = "Hey — what are you building?"

# Onboarding label → internal value maps
FOUNDER_TYPE_MAP = {
    "Student or aspiring founder": "student",
    "Working professional with a side project": "professional",
    "First-time founder": "founder",
    "Serial entrepreneur": "serial",
}

SECTOR_MAP = {
    "SaaS / B2B software": "saas",
    "D2C / Consumer brand": "d2c",
    "Fintech": "fintech",
    "Marketplace": "marketplace",
    "EdTech": "edtech",
    "HealthTech": "healthtech",
    "Climate / Deep Tech": "deeptech",
    "Something else": "unknown",
}

STAGE_MAP = {
    "Just an idea": "idea",
    "Researching and validating": "pre-revenue",
    "Actively building": "pre-revenue",
    "Have early users": "early-revenue",
    "Launched and growing": "growth",
}


# ── File parsing ──────────────────────────────────────────────────────────────

def _parse_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def _parse_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except ImportError:
        return "[python-docx not installed — run: pip install python-docx]"


def _parse_pdf(path: str) -> str:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    pages.append(f"[Page {i + 1}]\n{text.strip()}")
        return "\n\n".join(pages)
    except ImportError:
        return "[pdfplumber not installed — run: pip install pdfplumber]"


def _parse_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        slides = []
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        return "[python-pptx not installed — run: pip install python-pptx]"


def parse_uploaded_file(file_path: str) -> tuple[str, str]:
    """Parse an uploaded file. Returns (extracted_text, file_type_label)."""
    ext = Path(file_path).suffix.lower()
    if ext == ".txt":
        return _parse_txt(file_path), "notes"
    elif ext == ".docx":
        return _parse_docx(file_path), "document"
    elif ext == ".pdf":
        return _parse_pdf(file_path), "pitch deck"
    elif ext == ".pptx":
        return _parse_pptx(file_path), "pitch deck"
    return "", "file"


# ── LLM ──────────────────────────────────────────────────────────────────────

def _call_claude(messages: list[dict], system: str, max_tokens: int, model: str | None = None) -> str:
    """Call Claude via the Anthropic API."""
    import anthropic
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    response = client.messages.create(
        model=model or CLAUDE_MODEL_FAST,
        max_tokens=max_tokens,
        system=system,
        messages=messages,
    )
    return response.content[0].text


def _call_ollama(messages: list[dict], system: str, max_tokens: int) -> str:
    """Call local Ollama with the configured model."""
    import ollama as _ollama
    full_messages = [{"role": "system", "content": system}] + messages
    response = _ollama.chat(
        model=OLLAMA_MODEL,
        messages=full_messages,
        options={
            "num_predict": max_tokens,
            "temperature": 0.85,      # higher = more lexical variety
            "repeat_penalty": 1.18,   # penalise reusing the same words/phrases
            "top_p": 0.92,            # nucleus sampling — keeps output coherent but varied
            "top_k": 50,              # slightly wider vocabulary selection
        },
        think=False,
    )
    text = response.message.content or ""
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()
    return text


def call_llm(messages: list[dict], system: str, max_tokens: int = 1024, model_override: str | None = None) -> str:
    """Route to Claude API or local Ollama. model_override selects premium model for Crystallize."""
    if USE_CLAUDE:
        return _call_claude(messages, system, max_tokens, model=model_override or CLAUDE_MODEL_FAST)
    return _call_ollama(messages, system, max_tokens)


def analyze_state(conversation_history: list[dict], current_state: ConversationState) -> ConversationState:
    """Use the LLM to analyze the conversation and update state."""
    analysis_input = build_analysis_prompt(conversation_history)
    try:
        raw = call_llm(
            messages=[{"role": "user", "content": analysis_input}],
            system=ANALYSIS_SYSTEM_PROMPT,
            max_tokens=512,
        )
        cleaned = raw.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1].rsplit("```", 1)[0]
        start = cleaned.find("{")
        end = cleaned.rfind("}") + 1
        if start != -1 and end > start:
            cleaned = cleaned[start:end]
        current_state.update_from_analysis(json.loads(cleaned))
    except (json.JSONDecodeError, KeyError, IndexError, ValueError):
        current_state.turns += 1
    return current_state


# ── Display helpers ───────────────────────────────────────────────────────────

SHARPNESS_LABELS = {
    "Problem":        ("fuzzy", "emerging", "clear", "sharp"),
    "Solution":       ("vague", "forming", "defined", "polished"),
    "Market":         ("guesswork", "assumption-based", "researched", "data-backed"),
    "Business Model": ("unclear", "conceptual", "modeled", "unit-economics-ready"),
    "Traction":       ("none", "anecdotal", "early signals", "proven"),
    "Team":           ("unstated", "mentioned", "relevant", "compelling"),
    "Ask":            ("undecided", "rough", "calculated", "investor-ready"),
}

SECTION_ICONS = {
    "Problem": "🔍", "Solution": "💡", "Market": "📊",
    "Business Model": "💰", "Traction": "📈", "Team": "👥", "Ask": "🎯",
}


def _sharpness(section: str, score: int) -> str:
    labels = SHARPNESS_LABELS.get(section, ("low", "medium", "high", "complete"))
    if score < 25:
        return labels[0]
    elif score < 50:
        return labels[1]
    elif score < 75:
        return labels[2]
    return labels[3]


def format_coverage_display(state: ConversationState) -> str:
    """Format the coverage tracker for sidebar display."""
    lines = []

    for section, score in state.coverage.items():
        icon = SECTION_ICONS.get(section, "•")
        pct = score
        bar_filled = int(pct / 10)
        bar = "▓" * bar_filled + "░" * (10 - bar_filled)
        check = " ✓" if pct >= 70 else ""
        sharpness = _sharpness(section, pct)
        lines.append(f"{icon} **{section}**{check}  `{bar}`\n*{sharpness}*\n")

    overall = state.overall_coverage()

    if overall < 20:
        status = "Just getting started"
    elif overall < 45:
        status = "Building momentum"
    elif overall < 70:
        status = "Getting sharper"
    elif overall < 90:
        status = "Almost pitch-ready"
    else:
        status = "Pitch-ready 🚀"

    lines.append(f"---\n**{status}** — {overall:.0f}% explored")

    if state.sector != "unknown":
        lines.append(f"\n🏷 **Sector** · {state.sector.upper()}")
    if state.stage != "unknown":
        lines.append(f"📍 **Stage** · {state.stage.replace('-', ' ').title()}")
    if state.company_name:
        lines.append(f"🏢 **{state.company_name}**")
    if state.mode == "quick_stress_test":
        lines.append(f"\n⚡ **Quick Stress Test** mode")
    if state.urgency:
        lines.append(f"\n🔥 Urgency detected")

    unchallenged = state.get_unchallenged_numbers()
    if unchallenged:
        lines.append(f"\n⚠️ **{len(unchallenged)} assumption{'' if len(unchallenged) == 1 else 's'} to pressure-test**")

    return "\n".join(lines)


def deserialize_state(state_json: str) -> ConversationState:
    state = ConversationState()
    if state_json:
        d = json.loads(state_json)
        state.coverage = d.get("coverage", state.coverage)
        state.sector = d.get("sector", "unknown")
        state.stage = d.get("stage", "unknown")
        state.company_name = d.get("company_name", "")
        state.founder_type = d.get("founder_type", "unknown")
        state.mode = d.get("mode", "think_it_through")
        state.urgency = d.get("urgency", False)
        state.phase = d.get("phase", "intro")
        state.turns = d.get("turns", 0)
        state.facts = d.get("facts", {})
    return state


# ── Core chat logic ───────────────────────────────────────────────────────────

def chat(user_message: str, history: list, state_json: str, session_id: str = ""):
    state = deserialize_state(state_json)
    llm_messages = [{"role": m["role"], "content": m["content"]} for m in history]
    llm_messages.append({"role": "user", "content": user_message})

    # RAG: retrieve relevant knowledge + similar past conversations
    rag_context = ""
    if RAG_ENABLED:
        try:
            query = user_message + (" " + history[-1]["content"] if history else "")
            k_chunks = _rag.retrieve_knowledge(query, top_k=3, sector_filter=state.sector)
            c_chunks = _rag.retrieve_conversations(
                query, top_k=2, sector=state.sector, founder_type=state.founder_type
            )
            rag_context = _rag.format_rag_context(k_chunks, c_chunks)
        except Exception:
            pass

    system = build_system_prompt(
        state.to_json(),
        state.sector,
        state.phase,
        state.coverage.get("Ask", 0),
        state.mode,
        state.urgency,
        state.founder_type,
        rag_context=rag_context,
    )
    response = call_llm(llm_messages, system)
    history = history + [
        {"role": "user", "content": user_message},
        {"role": "assistant", "content": response},
    ]
    state = analyze_state(history, state)

    # Persist turn + update session
    if RAG_ENABLED and session_id:
        meta = {
            "founder_type": state.founder_type,
            "sector": state.sector,
            "stage": state.stage,
            "phase": state.phase,
        }
        try:
            _memory.store_turn(session_id, "user", user_message)
            _memory.store_turn(session_id, "assistant", response)
            _memory.update_session(session_id, state)
            _rag.store_conversation_turn(session_id, "user", user_message, meta)
            _rag.store_conversation_turn(session_id, "assistant", response, meta)
        except Exception:
            pass

    return "", history, state.to_json(), format_coverage_display(state)


def structure_pitch(history: list, state_json: str):
    if not history:
        return gr.update(
            value="**Start a conversation first** — share your idea and we'll explore it together before structuring.",
            visible=True,
        )
    state = deserialize_state(state_json)
    prompt = build_structure_prompt(state.to_json(), history)
    outline = call_llm(
        messages=[{"role": "user", "content": prompt}],
        system="You are a pitch deck structuring assistant. Create clear, actionable content outlines based on Socratic mentoring conversations.",
        max_tokens=2048,
        model_override=CLAUDE_MODEL_PREMIUM if USE_CLAUDE else None,
    )
    return gr.update(value=outline, visible=True)


# ── CSS ───────────────────────────────────────────────────────────────────────

CUSTOM_CSS = """
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

/* ── Base ── */
body, .gradio-container {
    font-family: 'Inter', system-ui, -apple-system, sans-serif !important;
    background: #0c0c10 !important;
    color: #e2e2f0 !important;
    min-height: 100vh;
}
.gradio-container .block, .gradio-container .form, .gradio-container .panel {
    background: transparent !important;
    border: none !important;
    box-shadow: none !important;
}

/* ── Slim top bar ── */
.topbar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0.7rem 1.5rem 0.5rem;
    border-bottom: 1px solid #1e1e2e;
    background: #0c0c10;
}
.topbar-brand {
    display: flex;
    align-items: center;
    gap: 0.6rem;
}
.topbar-brand .dot {
    width: 8px; height: 8px;
    border-radius: 50%;
    background: linear-gradient(135deg, #7c6af7, #38bdf8);
    box-shadow: 0 0 8px rgba(124,106,247,0.6);
}
.topbar-brand span {
    font-size: 0.95rem;
    font-weight: 600;
    letter-spacing: -0.2px;
    background: linear-gradient(90deg, #a78bfa, #38bdf8);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
}
.topbar-sub {
    font-size: 0.78rem;
    color: #4a4a6a;
    letter-spacing: 0.3px;
}

/* ── Onboarding panel ── */
.onboard-wrap {
    display: flex;
    justify-content: center;
    padding: 3rem 1rem 2rem;
    min-height: calc(100vh - 48px);
}
.onboard-inner {
    width: 100%;
    max-width: 560px;
}
.onboard-step-dots {
    display: flex;
    gap: 6px;
    margin-bottom: 2.5rem;
}
.onboard-dot {
    width: 6px; height: 6px;
    border-radius: 50%;
    background: #252540;
    transition: background 0.3s;
}
.onboard-dot.active { background: #7c6af7 !important; }

.onboard-title {
    font-size: 1.35rem;
    font-weight: 600;
    color: #d0d0f0;
    margin-bottom: 0.35rem;
    letter-spacing: -0.3px;
}
.onboard-sub {
    font-size: 0.83rem;
    color: #4a4a6a;
    margin-bottom: 1.75rem;
}

/* Radio cards */
.card-radio .wrap {
    display: flex !important;
    flex-direction: column !important;
    gap: 0.45rem !important;
    background: transparent !important;
    border: none !important;
    padding: 0 !important;
}
.card-radio label {
    background: #111120 !important;
    border: 1px solid #1e1e30 !important;
    border-radius: 10px !important;
    padding: 0.8rem 1.1rem !important;
    cursor: pointer !important;
    transition: all 0.18s !important;
    color: #7878a8 !important;
    font-size: 0.88rem !important;
    line-height: 1.4 !important;
}
.card-radio label:hover {
    border-color: #5549aa !important;
    background: #161628 !important;
    color: #c0c0de !important;
}
.card-radio label.selected,
.card-radio input[type=radio]:checked ~ label {
    border-color: #7c6af7 !important;
    background: #1a1a38 !important;
    color: #d8d0ff !important;
}
.card-radio > div { background: transparent !important; border: none !important; }
.card-radio .svelte-1gfkn6j { background: transparent !important; }

/* Skip link */
.skip-link {
    text-align: center;
    margin-top: 1.5rem;
    font-size: 0.8rem;
    color: #333355;
}
.skip-link button {
    background: transparent !important;
    border: none !important;
    color: #44446a !important;
    font-size: 0.8rem !important;
    cursor: pointer !important;
    padding: 0 !important;
    text-decoration: underline !important;
    text-underline-offset: 2px !important;
}
.skip-link button:hover { color: #7777aa !important; }

/* ── Main layout ── */
.main-row {
    display: flex !important;
    gap: 0 !important;
    padding: 0 !important;
    min-height: calc(100vh - 48px);
}

/* ── LEFT: Dashboard panel ── */
.left-panel {
    background: #0e0e16 !important;
    border-right: 1px solid #1e1e2e !important;
    padding: 1.25rem 1rem 1rem !important;
    display: flex;
    flex-direction: column;
    gap: 0.25rem;
}

.panel-label {
    font-size: 0.65rem;
    font-weight: 600;
    letter-spacing: 1.2px;
    text-transform: uppercase;
    color: #44446a;
    margin-bottom: 0.6rem;
    padding-bottom: 0.4rem;
    border-bottom: 1px solid #1a1a28;
}

/* Coverage markdown */
.left-panel .markdown-body, .left-panel .prose {
    font-size: 0.78rem !important;
    color: #9898bc !important;
    line-height: 1.5 !important;
}
.left-panel code {
    font-size: 0.7rem !important;
    letter-spacing: 1.5px !important;
    color: #6655cc !important;
    background: #15152a !important;
    padding: 1px 4px !important;
    border-radius: 3px !important;
}
.left-panel strong { color: #c0c0de !important; }
.left-panel em { color: #5555aa !important; font-style: normal !important; font-size: 0.72rem !important; }

/* Left panel buttons */
.left-panel button {
    width: 100% !important;
    border-radius: 8px !important;
    font-size: 0.82rem !important;
    font-weight: 500 !important;
}
.btn-crystallize {
    background: linear-gradient(135deg, #3730a3, #5b21b6) !important;
    border: none !important;
    color: #e0d8ff !important;
    padding: 0.6rem !important;
    margin-top: 0.75rem !important;
}
.btn-crystallize:hover {
    background: linear-gradient(135deg, #4338ca, #6d28d9) !important;
    transform: translateY(-1px) !important;
}
.btn-reset {
    background: transparent !important;
    border: 1px solid #1e1e30 !important;
    color: #44446a !important;
    padding: 0.4rem !important;
    margin-top: 0.3rem !important;
    font-size: 0.75rem !important;
}
.btn-reset:hover {
    border-color: #33334a !important;
    color: #7777aa !important;
}

/* Footnote */
.footnote {
    margin-top: auto;
    padding-top: 1rem;
    font-size: 0.72rem;
    color: #333355;
    line-height: 1.6;
    border-top: 1px solid #151520;
}
.footnote em { color: #44446a; font-style: normal; }

/* ── RIGHT: Chat area ── */
.right-panel {
    background: #0c0c10 !important;
    padding: 0.75rem 1.25rem 1rem !important;
    display: flex;
    flex-direction: column;
    gap: 0.5rem;
}

/* File upload */
.upload-strip {
    margin-bottom: 0.25rem !important;
}
.gradio-container .upload-strip label {
    font-size: 0.75rem !important;
    color: #44446a !important;
}
.gradio-container .upload-strip .file-preview,
.gradio-container .upload-strip [data-testid="upload-button"],
.gradio-container .upload-strip .wrap {
    background: #0e0e18 !important;
    border: 1px dashed #252540 !important;
    border-radius: 8px !important;
    color: #44446a !important;
    font-size: 0.78rem !important;
    transition: border-color 0.2s !important;
}
.gradio-container .upload-strip .wrap:hover {
    border-color: #5549aa !important;
}

/* Chat window */
.gradio-container .chatbot {
    background: #0c0c10 !important;
    border: none !important;
    border-radius: 0 !important;
    flex: 1;
}

/* User bubble */
.gradio-container .chatbot .message.user {
    background: #1a1a2e !important;
    color: #d0ccf0 !important;
    border: 1px solid #2a2a4a !important;
    border-radius: 16px 16px 4px 16px !important;
    font-size: 0.9rem !important;
    line-height: 1.65 !important;
    max-width: 80% !important;
    margin-left: auto !important;
}

/* Assistant bubble */
.gradio-container .chatbot .message.bot {
    background: transparent !important;
    color: #c8c8e8 !important;
    border: none !important;
    border-radius: 0 !important;
    font-size: 0.92rem !important;
    line-height: 1.75 !important;
    padding-left: 0 !important;
}

/* Chip buttons */
.chips-row {
    display: flex !important;
    flex-wrap: wrap !important;
    gap: 0.35rem !important;
    padding: 0.15rem 0 !important;
    min-height: 30px !important;
}
.chip-btn button {
    background: #0f0f1c !important;
    border: 1px solid #1e1e2e !important;
    border-radius: 20px !important;
    color: #5555aa !important;
    font-size: 0.77rem !important;
    padding: 0.2rem 0.75rem !important;
    height: auto !important;
    min-width: unset !important;
    white-space: nowrap !important;
    transition: all 0.15s !important;
    font-weight: 400 !important;
}
.chip-btn button:hover {
    border-color: #5549aa !important;
    color: #a0a0cc !important;
    background: #14142a !important;
}

/* Input row */
.input-row {
    margin-top: 0.25rem !important;
    align-items: flex-end !important;
}
.gradio-container .input-row textarea {
    background: #13131e !important;
    border: 1px solid #252540 !important;
    color: #e2e2f0 !important;
    border-radius: 12px !important;
    font-family: inherit !important;
    font-size: 0.9rem !important;
    line-height: 1.5 !important;
    resize: none !important;
    transition: border-color 0.2s !important;
}
.gradio-container .input-row textarea:focus {
    border-color: #5549aa !important;
    box-shadow: 0 0 0 3px rgba(85,73,170,0.12) !important;
    outline: none !important;
}
.gradio-container .input-row textarea::placeholder { color: #33334a !important; }

.gradio-container button.primary {
    background: linear-gradient(135deg, #4f46e5, #7c3aed) !important;
    border: none !important;
    border-radius: 10px !important;
    color: #fff !important;
    font-weight: 600 !important;
    font-size: 0.88rem !important;
    height: 52px !important;
    transition: opacity 0.15s, transform 0.1s !important;
}
.gradio-container button.primary:hover {
    opacity: 0.88 !important;
    transform: translateY(-1px) !important;
}

/* Structure output */
.structure-out {
    background: #0e0e18 !important;
    border: 1px solid #2a2a4a !important;
    border-radius: 12px !important;
    padding: 1.5rem !important;
    margin-top: 0.5rem !important;
    color: #c8c8e8 !important;
    font-size: 0.88rem !important;
    line-height: 1.75 !important;
}

/* Scrollbars */
::-webkit-scrollbar { width: 4px; height: 4px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: #1e1e30; border-radius: 2px; }
::-webkit-scrollbar-thumb:hover { background: #3a3a5e; }
"""


# ── UI ────────────────────────────────────────────────────────────────────────

def _dot_html(active_step: int) -> str:
    dots = ""
    for i in range(4):
        cls = "onboard-dot active" if i == active_step else "onboard-dot"
        dots += f'<span class="{cls}"></span>'
    return f'<div class="onboard-step-dots">{dots}</div>'


with gr.Blocks(title="Pitch Deck Mentor", css=CUSTOM_CSS, theme=gr.themes.Base()) as app:

    # ── Top bar ──
    gr.HTML("""
    <div class="topbar">
        <div class="topbar-brand">
            <div class="dot"></div>
            <span>Pitch Deck Mentor</span>
        </div>
        <div class="topbar-sub">Think like an investor. Build like a founder.</div>
    </div>
    """)

    state_json = gr.State("")
    conversation_history = gr.State([])
    onboard_data = gr.State({"founder_type": "", "sector": "", "stage": ""})
    session_id = gr.State("")

    # ── ONBOARDING PANEL ──────────────────────────────────────────────────────
    with gr.Column(visible=True, elem_classes=["onboard-wrap"]) as onboarding_col:
        with gr.Column(elem_classes=["onboard-inner"]):

            # Step 1: Founder type
            with gr.Column(visible=True) as step1_col:
                gr.HTML(_dot_html(0))
                gr.HTML('<div class="onboard-title">What best describes you right now?</div>'
                        '<div class="onboard-sub">This helps calibrate the depth and tone of the conversation.</div>')
                founder_type_radio = gr.Radio(
                    choices=list(FOUNDER_TYPE_MAP.keys()),
                    label="",
                    elem_classes=["card-radio"],
                )
                next_btn_1 = gr.Button("Continue →", variant="primary")
                with gr.Row(elem_classes=["skip-link"]):
                    skip_btn = gr.Button("Skip onboarding →")

            # Step 2: Sector
            with gr.Column(visible=False) as step2_col:
                gr.HTML(_dot_html(1))
                gr.HTML('<div class="onboard-title">What space is your idea in?</div>'
                        '<div class="onboard-sub">The right sector loads the right questions — D2C founders get challenged differently from SaaS founders.</div>')
                sector_radio = gr.Radio(
                    choices=list(SECTOR_MAP.keys()),
                    label="",
                    elem_classes=["card-radio"],
                )
                next_btn_2 = gr.Button("Continue →", variant="primary")

            # Step 3: Stage
            with gr.Column(visible=False) as step3_col:
                gr.HTML(_dot_html(2))
                gr.HTML('<div class="onboard-title">Where are you with this idea?</div>'
                        '<div class="onboard-sub">This is the most important calibration signal — it sets the depth of every question.</div>')
                stage_radio = gr.Radio(
                    choices=list(STAGE_MAP.keys()),
                    label="",
                    elem_classes=["card-radio"],
                )
                next_btn_3 = gr.Button("Continue →", variant="primary")

            # Step 4: Mode selection
            with gr.Column(visible=False) as step4_col:
                gr.HTML(_dot_html(3))
                gr.HTML('<div class="onboard-title">What do you need right now?</div>'
                        '<div class="onboard-sub">You can always ask to switch modes mid-conversation.</div>')
                mode_radio = gr.Radio(
                    choices=[
                        "Think It Through — explore your idea in full",
                        "Quick Stress Test — find the gaps fast",
                    ],
                    label="",
                    elem_classes=["card-radio"],
                )
                user_id_input = gr.Textbox(
                    placeholder="Optional — your name or email to resume sessions later",
                    show_label=False,
                    max_lines=1,
                )
                start_btn = gr.Button("Let's go →", variant="primary")

    # ── MAIN CHAT LAYOUT ──────────────────────────────────────────────────────
    with gr.Row(visible=False, elem_classes=["main-row"]) as main_row:

        # Left: Dashboard
        with gr.Column(scale=1, elem_classes=["left-panel"]):
            gr.HTML('<div class="panel-label">Pitch Coverage</div>')

            coverage_md = gr.Markdown(
                value=format_coverage_display(ConversationState()),
            )

            structure_btn = gr.Button(
                "✦  Crystallize My Pitch",
                variant="primary",
                elem_classes=["btn-crystallize"],
            )

            reset_btn = gr.Button(
                "↺  New conversation",
                variant="secondary",
                elem_classes=["btn-reset"],
            )

            gr.HTML("""
            <div class="footnote">
                A Socratic thinking partner for founders.<br/>
                <em>No templates. No boxes. Just the questions that matter for your specific business.</em>
            </div>
            """)

        # Right: Chat
        with gr.Column(scale=3, elem_classes=["right-panel"]):

            # File upload
            with gr.Row(elem_classes=["upload-strip"]):
                file_upload = gr.File(
                    label="📎  Drop a pitch deck, doc, or notes  ·  PDF  PPTX  DOCX  TXT",
                    file_types=[".pdf", ".pptx", ".docx", ".txt"],
                    file_count="single",
                    height=60,
                )

            # Chat window
            chatbot = gr.Chatbot(
                value=[],
                height=460,
                show_label=False,
                type="messages",
                avatar_images=(
                    None,
                    "https://api.dicebear.com/7.x/bottts-neutral/svg?seed=mentor&backgroundColor=4f46e5",
                ),
            )

            # Smart suggestion chips
            with gr.Row(elem_classes=["chips-row"]):
                chip1 = gr.Button("", visible=False, elem_classes=["chip-btn"])
                chip2 = gr.Button("", visible=False, elem_classes=["chip-btn"])
                chip3 = gr.Button("", visible=False, elem_classes=["chip-btn"])
                chip4 = gr.Button("", visible=False, elem_classes=["chip-btn"])

            # Input row
            with gr.Row(elem_classes=["input-row"]):
                msg = gr.Textbox(
                    placeholder="Type your answer or idea...",
                    show_label=False,
                    scale=5,
                    lines=2,
                    max_lines=5,
                )
                send_btn = gr.Button("Send", variant="primary", scale=1, min_width=80)

    structure_output = gr.Markdown(
        visible=False,
        elem_classes=["structure-out"],
    )

    # ── Event handlers ────────────────────────────────────────────────────────

    # All outputs that on_send and chip handlers need to return
    CHAT_OUTPUTS = [msg, file_upload, chatbot, conversation_history, state_json, coverage_md,
                    chip1, chip2, chip3, chip4]

    def _chip_updates_from_msg(assistant_msg: str, state: ConversationState):
        chips = generate_contextual_chips(assistant_msg, state)
        return [gr.update(value=c["text"], visible=c["visible"]) for c in chips]

    def _chip_updates_from_state(state: ConversationState):
        chips = get_chip_suggestions(state)
        return [gr.update(value=c["text"], visible=c["visible"]) for c in chips]

    def on_send(user_message, file_obj, displayed_chat, history, st_json, cov_display, sid=""):
        text = (user_message or "").strip()
        doc_content, file_type, file_name = "", "", ""

        if file_obj is not None:
            doc_content, file_type = parse_uploaded_file(file_obj.name)
            file_name = Path(file_obj.name).name

        if not text and not doc_content:
            state = deserialize_state(st_json)
            chips = _chip_updates_from_state(state)
            return ("", None, displayed_chat, history, st_json, cov_display, *chips)

        # LLM message includes full document text
        if doc_content and text:
            llm_message = f"{text}\n\n[Attached {file_type}: {file_name}]\n\n{doc_content}"
        elif doc_content:
            llm_message = f"[Shared {file_type}: {file_name}]\n\n{doc_content}"
        else:
            llm_message = text

        # Display message — no raw document dump
        if doc_content and text:
            display_message = f"{text}\n\n*📎 {file_name}*"
        elif doc_content:
            display_message = f"*📎 Shared {file_type}: {file_name}*"
        else:
            display_message = text

        displayed_chat = displayed_chat + [{"role": "user", "content": display_message}]

        _, updated_history, new_state, new_coverage = chat(llm_message, history, st_json, sid)

        assistant_msg = updated_history[-1]["content"]
        displayed_chat = displayed_chat + [{"role": "assistant", "content": assistant_msg}]

        updated_state_obj = deserialize_state(new_state)
        # Dynamic chips: derived from what the mentor just asked
        chips = _chip_updates_from_msg(assistant_msg, updated_state_obj)

        return ("", None, displayed_chat, updated_history, new_state, new_coverage, *chips)

    def on_chip_click(chip_val, file_obj, displayed_chat, history, st_json, cov_display, sid=""):
        if not chip_val:
            state = deserialize_state(st_json)
            chips = _chip_updates_from_state(state)
            return ("", file_obj, displayed_chat, history, st_json, cov_display, *chips)
        return on_send(chip_val, file_obj, displayed_chat, history, st_json, cov_display, sid)

    def on_structure(displayed_chat, history, st_json):
        if not history:
            return gr.update(
                value="**Start a conversation first** — tell me about your startup and we'll work through it together.",
                visible=True,
            )
        return structure_pitch(history, st_json)

    def on_reset():
        fresh = ConversationState()
        chips = _chip_updates_from_state(fresh)
        return (
            gr.update(visible=True),     # onboarding_col — go back to onboarding
            gr.update(visible=False),    # main_row
            gr.update(visible=True),     # step1_col
            gr.update(visible=False),    # step2_col
            gr.update(visible=False),    # step3_col
            gr.update(visible=False),    # step4_col
            [],                          # chatbot
            [],                          # conversation_history
            "",                          # state_json
            format_coverage_display(fresh),  # coverage_md
            gr.update(visible=False),    # structure_output
            {"founder_type": "", "sector": "", "stage": ""},  # onboard_data
            "",                          # session_id
            *chips,                      # chip1-4
        )

    # Onboarding step handlers
    def on_next_step1(selection, data):
        ft = selection or "First-time founder"
        new_data = {**data, "founder_type": ft}
        return gr.update(visible=False), gr.update(visible=True), new_data

    def on_next_step2(selection, data):
        sec = selection or "Something else"
        new_data = {**data, "sector": sec}
        return gr.update(visible=False), gr.update(visible=True), new_data

    def on_next_step3(selection, data):
        st = selection or "Just an idea"
        new_data = {**data, "stage": st}
        return gr.update(visible=False), gr.update(visible=True), new_data

    def on_start(mode_selection, data, user_id=""):
        """Complete onboarding and start the personalised conversation."""
        ft_label = data.get("founder_type", "") or "First-time founder"
        sec_label = data.get("sector", "") or "Something else"
        st_label = data.get("stage", "") or "Just an idea"

        state = ConversationState()
        state.founder_type = FOUNDER_TYPE_MAP.get(ft_label, "founder")
        state.sector = SECTOR_MAP.get(sec_label, "unknown")
        state.stage = STAGE_MAP.get(st_label, "idea")
        state.mode = (
            "quick_stress_test"
            if mode_selection and "Stress" in mode_selection
            else "think_it_through"
        )

        opening = build_personalized_opening(state.founder_type, state.sector, state.stage)
        initial_chat = [{"role": "assistant", "content": opening}]
        chips = _chip_updates_from_msg(opening, state)

        new_session_id = ""
        if RAG_ENABLED:
            try:
                new_session_id = _memory.create_session(state, user_id or "")
            except Exception:
                pass

        return (
            gr.update(visible=False),          # onboarding_col
            gr.update(visible=True),           # main_row
            initial_chat,                      # chatbot
            [],                                # conversation_history
            state.to_json(),                   # state_json
            format_coverage_display(state),    # coverage_md
            new_session_id,                    # session_id
            *chips,                            # chip1-4
        )

    def on_skip():
        """Skip onboarding and start with defaults."""
        state = ConversationState()
        initial_chat = [{"role": "assistant", "content": WELCOME_MESSAGE}]
        chips = _chip_updates_from_state(state)

        new_session_id = ""
        if RAG_ENABLED:
            try:
                new_session_id = _memory.create_session(state, "")
            except Exception:
                pass

        return (
            gr.update(visible=False),
            gr.update(visible=True),
            initial_chat,
            [],
            state.to_json(),
            format_coverage_display(state),
            new_session_id,
            *chips,
        )

    # Onboarding start outputs
    START_OUTPUTS = [onboarding_col, main_row, chatbot, conversation_history,
                     state_json, coverage_md, session_id, chip1, chip2, chip3, chip4]

    # Wire onboarding steps
    next_btn_1.click(on_next_step1, [founder_type_radio, onboard_data], [step1_col, step2_col, onboard_data])
    next_btn_2.click(on_next_step2, [sector_radio, onboard_data], [step2_col, step3_col, onboard_data])
    next_btn_3.click(on_next_step3, [stage_radio, onboard_data], [step3_col, step4_col, onboard_data])
    start_btn.click(on_start, [mode_radio, onboard_data, user_id_input], START_OUTPUTS)
    skip_btn.click(on_skip, [], START_OUTPUTS)

    # Chat handlers — io_base includes session_id for memory storage
    io_base = [msg, file_upload, chatbot, conversation_history, state_json, coverage_md, session_id]

    send_btn.click(on_send, io_base, CHAT_OUTPUTS)
    msg.submit(on_send, io_base, CHAT_OUTPUTS)

    # Chip click handlers
    for chip_btn in [chip1, chip2, chip3, chip4]:
        chip_btn.click(
            on_chip_click,
            [chip_btn] + io_base[1:],  # chip value + file_upload, chatbot, history, state, coverage, session_id
            CHAT_OUTPUTS,
        )

    # Structure + Reset
    structure_btn.click(on_structure, [chatbot, conversation_history, state_json], structure_output)

    reset_btn.click(
        on_reset,
        inputs=[],
        outputs=[
            onboarding_col, main_row,
            step1_col, step2_col, step3_col, step4_col,
            chatbot, conversation_history, state_json, coverage_md,
            structure_output, onboard_data, session_id,
            chip1, chip2, chip3, chip4,
        ],
    )


# ── ADMIN INTERFACE ───────────────────────────────────────────────────────────

if IS_ADMIN:
    ADMIN_CSS = """
    body, .gradio-container { font-family: 'Inter', system-ui, sans-serif !important; }
    """

    _ADMIN_SECTORS = ["general", "saas", "d2c", "fintech", "marketplace", "edtech", "healthtech", "deeptech"]

    with gr.Blocks(title="PDM Admin", css=ADMIN_CSS) as admin_app:
        gr.Markdown("# Pitch Deck Mentor — Admin\nManage knowledge base, review sessions, export training data.")

        with gr.Tabs():

            # ── Tab 1: Knowledge Base ──────────────────────────────────────
            with gr.Tab("📚 Knowledge Base"):
                with gr.Row():
                    with gr.Column(scale=1):
                        gr.Markdown("### Add content")
                        kb_file = gr.File(
                            label="Upload document",
                            file_types=[".pdf", ".docx", ".txt", ".pptx"],
                            file_count="single",
                        )
                        kb_url = gr.Textbox(
                            label="Or paste a URL",
                            placeholder="https://...",
                            max_lines=1,
                        )
                        kb_sector = gr.Dropdown(
                            choices=_ADMIN_SECTORS,
                            value="general",
                            label="Sector tag",
                        )
                        kb_ingest_btn = gr.Button("Index", variant="primary")
                        kb_inbox_btn = gr.Button("Index knowledge_inbox/ folder")
                        kb_status = gr.Markdown("")

                        gr.Markdown("### Remove content")
                        kb_delete_src = gr.Textbox(
                            label="Source to delete",
                            placeholder="filename.pdf or full URL",
                            max_lines=1,
                        )
                        kb_delete_btn = gr.Button("Delete source", variant="stop")

                    with gr.Column(scale=2):
                        gr.Markdown("### Indexed sources")
                        kb_table = gr.Dataframe(
                            headers=["source", "doc_type", "sector", "chunk_count", "date_added"],
                            label="",
                            interactive=False,
                        )
                        kb_refresh_btn = gr.Button("Refresh")

                def _admin_ingest(file_obj, url, sector):
                    if not RAG_ENABLED:
                        return "⚠️ RAG not available — install chromadb and sentence-transformers."
                    msgs = []
                    if file_obj is not None:
                        try:
                            n = _rag.ingest_file(file_obj.name, sector=sector)
                            msgs.append(f"✅ Indexed **{Path(file_obj.name).name}** — {n} chunks")
                        except Exception as e:
                            msgs.append(f"❌ File error: {e}")
                    if url and url.strip():
                        try:
                            n = _rag.ingest_url(url.strip(), sector=sector)
                            msgs.append(f"✅ Indexed URL — {n} chunks")
                        except Exception as e:
                            msgs.append(f"❌ URL error: {e}")
                    return "\n\n".join(msgs) if msgs else "Nothing to index — provide a file or URL."

                def _admin_inbox():
                    if not RAG_ENABLED:
                        return "⚠️ RAG not available."
                    try:
                        n = _rag.ingest_inbox()
                        return f"✅ Indexed inbox — {n} new chunks"
                    except Exception as e:
                        return f"❌ Error: {e}"

                def _admin_refresh_sources():
                    if not RAG_ENABLED:
                        return []
                    sources = _rag.list_knowledge_sources()
                    return [[s["source"], s["doc_type"], s["sector"], s["chunk_count"], s["date_added"][:10]] for s in sources]

                def _admin_delete(source):
                    if not RAG_ENABLED:
                        return "⚠️ RAG not available."
                    if not source or not source.strip():
                        return "Enter a source name to delete."
                    try:
                        n = _rag.delete_source(source.strip())
                        return f"✅ Deleted {n} chunks for **{source.strip()}**" if n else f"⚠️ Source not found: {source.strip()}"
                    except Exception as e:
                        return f"❌ Error: {e}"

                kb_ingest_btn.click(_admin_ingest, [kb_file, kb_url, kb_sector], kb_status)
                kb_inbox_btn.click(_admin_inbox, [], kb_status)
                kb_refresh_btn.click(_admin_refresh_sources, [], kb_table)
                kb_delete_btn.click(_admin_delete, [kb_delete_src], kb_status)

            # ── Tab 2: Sessions ────────────────────────────────────────────
            with gr.Tab("👥 Sessions"):
                ses_refresh_btn = gr.Button("Refresh sessions")
                ses_table = gr.Dataframe(
                    headers=["id", "user", "sector", "stage", "founder_type", "company", "turns", "last_active"],
                    label="Recent sessions",
                    interactive=False,
                )
                with gr.Row():
                    ses_id_input = gr.Textbox(
                        label="Session ID",
                        placeholder="Paste a session UUID to inspect",
                        max_lines=1,
                    )
                    ses_view_btn = gr.Button("View conversation")
                ses_chat = gr.Chatbot(label="Conversation", type="messages", height=400)

                def _admin_list_sessions():
                    if not RAG_ENABLED:
                        return []
                    sessions = _memory.list_sessions(limit=200)
                    return [
                        [
                            s["id"][:8] + "…",
                            s.get("user_identifier", ""),
                            s.get("sector", ""),
                            s.get("stage", ""),
                            s.get("founder_type", ""),
                            s.get("company_name", ""),
                            s.get("turn_count", 0),
                            (s.get("last_active", "") or "")[:16],
                        ]
                        for s in sessions
                    ]

                def _admin_view_session(session_id_val):
                    if not RAG_ENABLED or not session_id_val.strip():
                        return []
                    try:
                        turns = _memory.get_session_turns(session_id_val.strip())
                        return [{"role": t["role"], "content": t["content"]} for t in turns if t["role"] in ("user", "assistant")]
                    except Exception as e:
                        return [{"role": "assistant", "content": f"Error: {e}"}]

                ses_refresh_btn.click(_admin_list_sessions, [], ses_table)
                ses_view_btn.click(_admin_view_session, [ses_id_input], ses_chat)

            # ── Tab 3: Export & Stats ──────────────────────────────────────
            with gr.Tab("📊 Export & Stats"):
                with gr.Row():
                    stats_refresh_btn = gr.Button("Refresh stats")
                    export_btn = gr.Button("Export fine-tuning JSONL", variant="primary")
                export_status = gr.Markdown("")
                stats_display = gr.JSON(label="Stats")

                def _admin_get_stats():
                    if not RAG_ENABLED:
                        return {"error": "RAG/Memory not available"}
                    stats = _memory.get_stats()
                    rag_stats = _rag.get_rag_stats()
                    return {**stats, "rag": rag_stats}

                def _admin_export():
                    if not RAG_ENABLED:
                        return "⚠️ Memory not available — install chromadb and sentence-transformers."
                    try:
                        path, count = _memory.export_jsonl()
                        return f"✅ Exported **{count} sessions** → `{path}`"
                    except Exception as e:
                        return f"❌ Export failed: {e}"

                stats_refresh_btn.click(_admin_get_stats, [], stats_display)
                export_btn.click(_admin_export, [], export_status)


# ── Launch ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if IS_ADMIN:
        if not RAG_ENABLED:
            print("⚠️  WARNING: chromadb / sentence-transformers not installed.")
            print("   Run: pip install chromadb sentence-transformers beautifulsoup4 requests")
            print("   Admin will launch but knowledge-base features will be disabled.\n")
        print("Admin interface → http://localhost:7861")
        admin_app.launch(server_port=7861, server_name="0.0.0.0", share=False)
    else:
        print("User interface → http://localhost:7860")
        app.launch(server_port=7860, server_name="0.0.0.0", share=False)

"""Upload parsing and lightweight lexical retrieval."""

from __future__ import annotations

import json
import mimetypes
import os
import re
import shutil
import subprocess
from datetime import datetime, timezone
from functools import lru_cache
from pathlib import Path

from fastapi import UploadFile


ROOT_DIR = Path(__file__).resolve().parents[2]
DATA_DIR = Path(os.environ.get("SIFT_DATA_DIR", str(ROOT_DIR / "data")))
UPLOAD_BACKEND = os.environ.get("SIFT_UPLOAD_BACKEND", "local").strip().lower()
_UPLOAD_TMP_VALUE = os.environ.get("SIFT_UPLOAD_TMP_DIR", "").strip()
UPLOAD_TMP_DIR = Path(_UPLOAD_TMP_VALUE or ("/tmp/sift_uploads" if UPLOAD_BACKEND == "gcs" else str(DATA_DIR / "session_uploads")))
UPLOADS_DIR = DATA_DIR / "session_uploads"
GCS_UPLOAD_BUCKET = os.environ.get("SIFT_UPLOAD_BUCKET", "").strip()
GCS_UPLOAD_PREFIX = os.environ.get("SIFT_UPLOAD_PREFIX", "session_uploads").strip().strip("/")
CHUNK_SIZE = 900
CHUNK_OVERLAP = 150
TOKEN_PATTERN = re.compile(r"[a-z0-9]{3,}")
ALLOWED_UPLOAD_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt"}
MAX_UPLOAD_BYTES = int(float(os.environ.get("SIFT_MAX_UPLOAD_MB", "20")) * 1024 * 1024)
MAX_UPLOADS_PER_SESSION = int(os.environ.get("SIFT_MAX_UPLOADS_PER_SESSION", "8"))


def _use_gcs() -> bool:
    return UPLOAD_BACKEND in {"gcs", "cloud_storage", "google_cloud_storage"}


def _project_id() -> str:
    return (
        os.environ.get("SIFT_GCP_PROJECT_ID")
        or os.environ.get("GOOGLE_CLOUD_PROJECT")
        or os.environ.get("GCP_PROJECT")
        or ""
    ).strip()


@lru_cache(maxsize=1)
def _gcs_bucket():
    if not GCS_UPLOAD_BUCKET:
        raise RuntimeError("SIFT_UPLOAD_BACKEND=gcs requires SIFT_UPLOAD_BUCKET.")
    try:
        from google.cloud import storage  # type: ignore
    except ImportError as exc:
        raise RuntimeError(
            "SIFT_UPLOAD_BACKEND=gcs requires google-cloud-storage. "
            "Install production requirements before starting the service."
        ) from exc

    project = _project_id() or None
    return storage.Client(project=project).bucket(GCS_UPLOAD_BUCKET)


def _gcs_object_name(session_id: str, filename: str) -> str:
    bits = [GCS_UPLOAD_PREFIX, session_id, filename] if GCS_UPLOAD_PREFIX else [session_id, filename]
    return "/".join(bit.strip("/") for bit in bits if bit)


def _upload_to_gcs(path: Path, session_id: str, filename: str) -> str:
    object_name = _gcs_object_name(session_id, filename)
    blob = _gcs_bucket().blob(object_name)
    content_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    blob.upload_from_filename(str(path), content_type=content_type)
    return object_name


def _load_gcs_json(session_id: str, filename: str, default):
    blob = _gcs_bucket().blob(_gcs_object_name(session_id, filename))
    try:
        if not blob.exists():
            return default
        return json.loads(blob.download_as_text())
    except json.JSONDecodeError:
        return default


def _download_gcs_object(object_name: str, destination: Path) -> Path:
    destination.parent.mkdir(parents=True, exist_ok=True)
    if destination.exists():
        return destination
    _gcs_bucket().blob(object_name).download_to_filename(str(destination))
    return destination


def _session_dir(session_id: str) -> Path:
    root = UPLOAD_TMP_DIR if _use_gcs() else UPLOADS_DIR
    path = root / session_id
    path.mkdir(parents=True, exist_ok=True)
    return path


def _manifest_path(session_id: str) -> Path:
    return _session_dir(session_id) / "manifest.json"


def _artifact_dir(session_id: str, stored_name: str) -> Path:
    path = _session_dir(session_id) / f"{stored_name}.artifacts"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _load_manifest(session_id: str) -> list[dict]:
    if _use_gcs():
        return _load_gcs_json(session_id, "manifest.json", [])

    path = _manifest_path(session_id)
    if not path.exists():
        return []
    try:
        return json.loads(path.read_text())
    except json.JSONDecodeError:
        return []


def _save_manifest(session_id: str, manifest: list[dict]) -> None:
    _manifest_path(session_id).write_text(json.dumps(manifest, indent=2))
    if _use_gcs():
        _upload_to_gcs(_manifest_path(session_id), session_id, "manifest.json")


def _sanitize_filename(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", name or "upload")
    return cleaned.strip("-") or "upload"


def _format_mb(byte_count: int) -> str:
    return f"{byte_count / (1024 * 1024):.0f}MB"


def _validate_upload(filename: str, content: bytes, existing_count: int) -> None:
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_UPLOAD_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_UPLOAD_EXTENSIONS))
        raise ValueError(f"Unsupported file type. Upload one of: {allowed}.")
    if not content:
        raise ValueError("That file is empty. Upload a PDF, PPTX, DOCX, or TXT with readable content.")
    if len(content) > MAX_UPLOAD_BYTES:
        raise ValueError(f"File is too large. Current beta limit is {_format_mb(MAX_UPLOAD_BYTES)} per upload.")
    if existing_count >= MAX_UPLOADS_PER_SESSION:
        raise ValueError(f"This session already has {MAX_UPLOADS_PER_SESSION} uploads. Start a new session to add more files.")


def _parse_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore").strip()


def _parse_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def _parse_pdf(path: Path) -> str:
    import pdfplumber

    pages = []
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {index + 1}]\n{text.strip()}")
    return "\n\n".join(pages)


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    slides = []
    prs = Presentation(path)
    for index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides.append(f"[Slide {index + 1}]\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def _slide_summary(text: str, limit: int = 240) -> str:
    cleaned = re.sub(r"\s+", " ", (text or "").strip())
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 3].rstrip() + "..."


def _render_pdf_page_images(path: Path, output_dir: Path) -> list[str]:
    renderer = shutil.which("pdftoppm")
    if not renderer:
        return []
    prefix = output_dir / "page"
    try:
        subprocess.run(
            [renderer, "-png", "-r", "110", str(path), str(prefix)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except Exception:
        return []
    pages = sorted(output_dir.glob("page-*.png"), key=lambda item: int(re.search(r"-(\d+)\.png$", item.name).group(1)) if re.search(r"-(\d+)\.png$", item.name) else 0)
    return [str(page) for page in pages]


def _build_pdf_deck_artifact(path: Path, session_id: str, stored_name: str) -> dict:
    import pdfplumber

    asset_dir = _artifact_dir(session_id, stored_name)
    image_paths = _render_pdf_page_images(path, asset_dir)
    slides = []
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages):
            text = (page.extract_text() or "").strip()
            image_path = image_paths[index] if index < len(image_paths) else ""
            slides.append(
                {
                    "index": index + 1,
                    "label": f"Page {index + 1}",
                    "extractedText": text,
                    "summary": _slide_summary(text or f"Page {index + 1}"),
                    "imagePath": image_path,
                }
            )
    limitations = []
    if not any(slide.get("imagePath") for slide in slides):
        limitations.append("Page images were not rendered in this environment, so review falls back to extracted text.")
    return {
        "artifactType": "deck",
        "format": "pdf",
        "storedName": stored_name,
        "slideCount": len(slides),
        "hasRenderableSlides": any(slide.get("imagePath") for slide in slides),
        "limitations": limitations,
        "slides": slides,
    }


def _build_pptx_deck_artifact(path: Path, stored_name: str) -> dict:
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        text = "\n".join(slide_text).strip()
        slides.append(
            {
                "index": index + 1,
                "label": f"Slide {index + 1}",
                "extractedText": text,
                "summary": _slide_summary(text or f"Slide {index + 1}"),
                "imagePath": "",
            }
        )
    return {
        "artifactType": "deck",
        "format": "pptx",
        "storedName": stored_name,
        "slideCount": len(slides),
        "hasRenderableSlides": False,
        "limitations": [
            "This environment extracted slide text from the PPTX but did not render slide images, so visual design claims should be treated as unverified.",
        ],
        "slides": slides,
    }


def _build_upload_artifact(path: Path, session_id: str, stored_name: str, doc_type: str) -> dict | None:
    ext = path.suffix.lower()
    if doc_type != "pitch deck":
        return None
    if ext == ".pdf":
        return _build_pdf_deck_artifact(path, session_id, stored_name)
    if ext == ".pptx":
        return _build_pptx_deck_artifact(path, stored_name)
    return None


def _persist_artifact_assets(session_id: str, stored_name: str, artifact: dict) -> None:
    if not _use_gcs():
        return
    for slide in artifact.get("slides", []):
        image_path_value = slide.get("imagePath", "") or ""
        if not image_path_value:
            continue
        image_path = Path(image_path_value)
        if not image_path.is_file():
            continue
        relative_name = f"{stored_name}.artifacts/{image_path.name}"
        slide["imageObject"] = _upload_to_gcs(image_path, session_id, relative_name)


def _materialize_artifact_assets(session_id: str, artifact: dict) -> dict:
    if not _use_gcs():
        return artifact
    stored_name = artifact.get("storedName", "deck")
    asset_dir = _artifact_dir(session_id, stored_name)
    for slide in artifact.get("slides", []):
        image_object = slide.get("imageObject", "")
        if not image_object:
            continue
        current_path_value = slide.get("imagePath", "") or ""
        current_path = Path(current_path_value) if current_path_value else None
        if current_path is not None and current_path.is_file():
            continue
        local_path = asset_dir / Path(image_object).name
        try:
            _download_gcs_object(image_object, local_path)
            slide["imagePath"] = str(local_path)
        except Exception:
            slide["imagePath"] = ""
    return artifact


def parse_uploaded_path(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    if ext == ".txt":
        return _parse_txt(path), "notes"
    if ext == ".docx":
        return _parse_docx(path), "document"
    if ext == ".pdf":
        return _parse_pdf(path), "pitch deck"
    if ext == ".pptx":
        return _parse_pptx(path), "pitch deck"
    return "", "file"


def _chunk_text(text: str, source: str, doc_type: str) -> list[dict]:
    chunks = []
    text = text.strip()
    cursor = 0
    index = 0
    while cursor < len(text):
        chunk = text[cursor: cursor + CHUNK_SIZE].strip()
        if chunk:
            chunks.append(
                {
                    "id": f"{source}:{index}",
                    "source": source,
                    "doc_type": doc_type,
                    "text": chunk,
                }
            )
        index += 1
        cursor += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


async def ingest_upload(session_id: str, upload: UploadFile) -> dict:
    content = await upload.read()
    filename = _sanitize_filename(upload.filename or "upload")
    manifest = _load_manifest(session_id)
    _validate_upload(filename, content, len(manifest))
    session_dir = _session_dir(session_id)
    stored_name = f"{datetime.now(timezone.utc).strftime('%Y%m%d%H%M%S')}-{filename}"
    stored_path = session_dir / stored_name
    stored_path.write_bytes(content)
    if _use_gcs():
        _upload_to_gcs(stored_path, session_id, stored_name)

    text, doc_type = parse_uploaded_path(stored_path)
    chunks = _chunk_text(text, filename, doc_type)
    chunks_path = session_dir / f"{stored_name}.chunks.json"
    chunks_path.write_text(json.dumps(chunks, indent=2))
    if _use_gcs():
        _upload_to_gcs(chunks_path, session_id, chunks_path.name)
    artifact = _build_upload_artifact(stored_path, session_id, stored_name, doc_type)
    artifact_path = None
    if artifact is not None:
        _persist_artifact_assets(session_id, stored_name, artifact)
        artifact_path = session_dir / f"{stored_name}.artifact.json"
        artifact_path.write_text(json.dumps(artifact, indent=2))
        if _use_gcs():
            _upload_to_gcs(artifact_path, session_id, artifact_path.name)

    entry = {
        "name": filename,
        "stored_name": stored_name,
        "docType": doc_type,
        "chunkCount": len(chunks),
        "chars": len(text),
        "bytes": len(content),
        "uploadedAt": datetime.now(timezone.utc).isoformat(),
        "chunksFile": chunks_path.name,
    }
    if artifact_path is not None:
        entry["artifactFile"] = artifact_path.name
        entry["slideCount"] = int(artifact.get("slideCount", 0) or 0)
        entry["hasRenderableSlides"] = bool(artifact.get("hasRenderableSlides", False))
    manifest.append(entry)
    _save_manifest(session_id, manifest)
    return entry


def list_active_uploads(session_id: str) -> list[dict]:
    public_fields = ("name", "docType", "chunkCount", "chars", "bytes", "uploadedAt", "slideCount", "hasRenderableSlides")
    return [
        {field: entry[field] for field in public_fields if field in entry}
        for entry in _load_manifest(session_id)
    ]


def load_deck_artifact(session_id: str, *, latest_only: bool = True) -> dict | None:
    manifest = _load_manifest(session_id)
    if not manifest:
        return None
    candidates = [entry for entry in manifest if entry.get("docType") == "pitch deck" and entry.get("artifactFile")]
    if not candidates:
        return None
    entries = [candidates[-1]] if latest_only else list(reversed(candidates))
    for entry in entries:
        if _use_gcs():
            artifact = _load_gcs_json(session_id, entry["artifactFile"], None)
            if not artifact:
                continue
            artifact = _materialize_artifact_assets(session_id, artifact)
        else:
            path = _session_dir(session_id) / entry["artifactFile"]
            if not path.exists():
                continue
            try:
                artifact = json.loads(path.read_text())
            except json.JSONDecodeError:
                continue
        artifact["name"] = entry.get("name", artifact.get("storedName", "deck"))
        artifact["uploadedAt"] = entry.get("uploadedAt", "")
        return artifact
    return None


def _load_chunks(session_id: str, chunks_file: str) -> list[dict]:
    if _use_gcs():
        return _load_gcs_json(session_id, chunks_file, [])
    chunks_path = _session_dir(session_id) / chunks_file
    if not chunks_path.exists():
        return []
    try:
        return json.loads(chunks_path.read_text())
    except json.JSONDecodeError:
        return []


def _tokenize(text: str) -> set[str]:
    return set(TOKEN_PATTERN.findall((text or "").lower()))


def retrieve_upload_context(
    session_id: str,
    query: str,
    top_k: int = 2,
    max_chars: int = 1200,
) -> list[dict]:
    manifest = _load_manifest(session_id)
    if not manifest:
        return []

    query_terms = _tokenize(query)
    candidates = []
    for entry in manifest:
        chunks = _load_chunks(session_id, entry["chunksFile"])
        if not chunks:
            continue
        for chunk in chunks:
            text = chunk.get("text", "")
            lower = text.lower()
            chunk_terms = _tokenize(text)
            overlap = len(query_terms & chunk_terms) if query_terms else 0
            phrase_bonus = 1 if any(term in lower for term in query_terms) else 0
            candidates.append(
                {
                    **chunk,
                    "score": overlap + phrase_bonus,
                    "uploadedAt": entry["uploadedAt"],
                }
            )

    candidates.sort(
        key=lambda item: (item["score"], item["uploadedAt"], -len(item["text"])),
        reverse=True,
    )
    selected = []
    total_chars = 0
    for chunk in candidates[: max(top_k * 4, top_k)]:
        snippet = chunk["text"].strip()
        if not snippet:
            continue
        if total_chars + len(snippet) > max_chars and selected:
            break
        selected.append(
            {
                "source": chunk["source"],
                "docType": chunk["doc_type"],
                "text": snippet[: min(len(snippet), 550)],
            }
        )
        total_chars += len(snippet)
        if len(selected) >= top_k:
            break

    if not selected and manifest:
        latest = manifest[-1]
        chunks = _load_chunks(session_id, latest["chunksFile"])
        selected = [
            {
                "source": latest["name"],
                "docType": latest["docType"],
                "text": chunk["text"][:550],
            }
            for chunk in chunks[:top_k]
        ]
    return selected


def upload_storage_status() -> dict[str, str | bool]:
    return {
        "backend": "gcs" if _use_gcs() else "local",
        "bucket": GCS_UPLOAD_BUCKET if _use_gcs() else "",
        "prefix": GCS_UPLOAD_PREFIX if _use_gcs() else "",
        "tmpDir": str(UPLOAD_TMP_DIR),
    }
